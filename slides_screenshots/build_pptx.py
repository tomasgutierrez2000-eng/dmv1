#!/usr/bin/env python3
"""
Build a professional DSCR Lineage Overview PowerPoint deck.
Clean minimal light theme with screenshots from the live application.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Constants ───────────────────────────────────────────────────────────────
SCREENSHOTS = '/Users/tomas/120/slides_screenshots'
OUTPUT = '/Users/tomas/120/DSCR_Lineage_Overview.pptx'

# Colors
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
DARK_GRAY = RGBColor(0x33, 0x33, 0x44)
MED_GRAY = RGBColor(0x66, 0x66, 0x77)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF5)
ORANGE = RGBColor(0xE8, 0x6C, 0x00)  # PwC-style orange
BLUE = RGBColor(0x3B, 0x82, 0xF6)
PURPLE = RGBColor(0xA7, 0x55, 0xF7)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
RED = RGBColor(0xEF, 0x44, 0x44)
TEAL = RGBColor(0x06, 0xB6, 0xD4)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


# ─── Helper Functions ────────────────────────────────────────────────────────

def add_background(slide, color=WHITE):
    """Set solid background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    """Add a filled rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    """Add a rounded rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 font_color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Calibri', line_spacing=1.15):
    """Add a text box with single-run text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.name = font_name
    if line_spacing != 1.0:
        p.line_spacing = line_spacing
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, font_size=14,
                       font_color=DARK_GRAY, font_name='Calibri', line_spacing=1.3):
    """Add text box with multiple paragraphs (each item is (text, bold, color))."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, bold, color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(4)
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = color if color else font_color
        run.font.bold = bold
        run.font.name = font_name
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=13,
                    font_color=DARK_GRAY, bullet_char='\u2022'):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        p.space_before = Pt(2)
        p.line_spacing = 1.3
        run = p.add_run()
        run.text = f'{bullet_char}  {item}'
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        run.font.name = 'Calibri'
    return txBox


def add_screenshot(slide, img_path, left, top, width, height=None, with_border=False):
    """Add a screenshot image, optionally constrained to width."""
    if os.path.exists(img_path):
        if height:
            pic = slide.shapes.add_picture(img_path, left, top, width, height)
        else:
            pic = slide.shapes.add_picture(img_path, left, top, width=width)
        if with_border:
            pic.line.color.rgb = RGBColor(0xCC, 0xCC, 0xDD)
            pic.line.width = Pt(1)
        return pic
    return None


def add_step_badge(slide, left, top, step_num, label, color=ORANGE):
    """Add a colored step badge."""
    badge = add_rounded_rect(slide, left, top, Inches(1.1), Inches(0.4), color)
    badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = badge.text_frame.paragraphs[0].add_run()
    run.text = f'STEP {step_num}'
    run.font.size = Pt(11)
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.name = 'Calibri'

    add_text_box(slide, left + Inches(1.25), top + Inches(0.02),
                 Inches(4), Inches(0.4), label, font_size=16,
                 font_color=DARK_GRAY, bold=True)


def add_accent_bar(slide, left, top, width, color=ORANGE):
    """Add a thin accent bar."""
    add_shape_rect(slide, left, top, width, Pt(3), color)


def add_footer(slide, slide_num, total=14):
    """Add consistent footer with page number."""
    # Bottom accent line
    add_shape_rect(slide, Inches(0), SLIDE_HEIGHT - Inches(0.4),
                   SLIDE_WIDTH, Pt(2), ORANGE)
    # Page number
    add_text_box(slide, SLIDE_WIDTH - Inches(1.5), SLIDE_HEIGHT - Inches(0.38),
                 Inches(1.3), Inches(0.3), f'{slide_num} / {total}',
                 font_size=9, font_color=MED_GRAY, alignment=PP_ALIGN.RIGHT)
    # Branding
    add_text_box(slide, Inches(0.5), SLIDE_HEIGHT - Inches(0.38),
                 Inches(3), Inches(0.3), 'DSCR Lineage Overview',
                 font_size=9, font_color=MED_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE SLIDE
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_background(slide, WHITE)

# Left color block
add_shape_rect(slide, Inches(0), Inches(0), Inches(0.4), SLIDE_HEIGHT, ORANGE)

# Title area
add_text_box(slide, Inches(1.2), Inches(1.8), Inches(7), Inches(1),
             'DSCR End-to-End Lineage', font_size=40, font_color=BLACK, bold=True)
add_accent_bar(slide, Inches(1.2), Inches(2.9), Inches(2.5), ORANGE)
add_text_box(slide, Inches(1.2), Inches(3.2), Inches(7), Inches(0.8),
             'From metric definition to dashboard\nHow DSCR is calculated, stored, and rolled up across the organization',
             font_size=18, font_color=MED_GRAY)

# Right side: screenshot preview
add_screenshot(slide, f'{SCREENSHOTS}/01_header_step1.png',
               Inches(7.5), Inches(1.2), Inches(5.3))

# Subtitle info
add_bullet_list(slide, Inches(1.2), Inches(4.5), Inches(5.5), Inches(2.5), [
    'What is DSCR and why it matters',
    'Two product variants: CRE (Real Estate) and C&I (Corporate)',
    'Where the data comes from (source tables)',
    'How the formula works at facility level',
    'The golden rule: never average pre-computed ratios',
    'How results roll up: Facility \u2192 Counterparty \u2192 Desk \u2192 Portfolio \u2192 LoB',
], font_size=14, font_color=DARK_GRAY)

add_footer(slide, 1)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: AGENDA / OVERVIEW
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
             'The 6-Step Lineage Flow', font_size=32, font_color=BLACK, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(2), ORANGE)
add_text_box(slide, Inches(0.8), Inches(1.35), Inches(10), Inches(0.5),
             'Every DSCR value on a dashboard traces back through these six stages',
             font_size=15, font_color=MED_GRAY)

# 6 step cards in 2 rows of 3
steps = [
    ('1', 'Metric Configuration', 'User defines which product variant to calculate\n(CRE or C&I) and selects the formula components.', BLUE),
    ('2', 'Dimensional Anchors', 'L1 reference tables identify the "who" and "what":\nfacilities, counterparties, products, and hierarchies.', TEAL),
    ('3', 'Source Data Tables', 'L2 snapshot tables hold the actual financial data:\nrevenue, expenses, debt payments, cash flows.', ORANGE),
    ('4', 'Facility Calculation', 'The DSCR formula is applied at facility level.\nNumerator \u00f7 Denominator = DSCR ratio.', GREEN),
    ('5', 'Storage & Rollup', 'Results stored in L3 tables, then rolled up through\n5 levels using the "pool and divide" method.', PURPLE),
    ('6', 'Dashboard Connection', 'Self-service: pick a variant and dimension level,\nand the platform auto-resolves the join logic.', RGBColor(0xEC, 0x48, 0x99)),
]

for i, (num, title, desc, color) in enumerate(steps):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + col * Inches(4.1)
    y = Inches(2.2) + row * Inches(2.5)

    card = add_rounded_rect(slide, x, y, Inches(3.8), Inches(2.1), LIGHT_GRAY, color)

    # Step number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.15), y + Inches(0.15),
                                     Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = circle.text_frame.paragraphs[0].add_run()
    run.text = num
    run.font.size = Pt(18)
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.name = 'Calibri'

    add_text_box(slide, x + Inches(0.8), y + Inches(0.18), Inches(2.8), Inches(0.4),
                 title, font_size=16, font_color=DARK_GRAY, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.75), Inches(3.4), Inches(1.2),
                 desc, font_size=12, font_color=MED_GRAY, line_spacing=1.3)

add_footer(slide, 2)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: WHAT IS DSCR?
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
             'What is DSCR?', font_size=32, font_color=BLACK, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(2), ORANGE)

# Main definition
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.7),
             'Debt Service Coverage Ratio measures how much cash flow a borrower generates '
             'compared to what they owe in loan payments.',
             font_size=17, font_color=DARK_GRAY)

# Formula box
formula_box = add_rounded_rect(slide, Inches(2.5), Inches(2.5), Inches(8), Inches(1.2),
                                RGBColor(0xF8, 0xF8, 0xFF), BLUE)
add_text_box(slide, Inches(3), Inches(2.6), Inches(7), Inches(1),
             'DSCR  =  Cash Flow Available for Debt Service  \u00f7  Total Debt Service',
             font_size=22, font_color=BLUE, bold=True, alignment=PP_ALIGN.CENTER,
             font_name='Consolas')

# What it means
add_text_box(slide, Inches(0.8), Inches(4.1), Inches(5.5), Inches(0.4),
             'What does the number mean?', font_size=16, font_color=DARK_GRAY, bold=True)

meanings = [
    'DSCR = 1.0x \u2192 Borrower earns exactly enough to cover debt payments (break-even)',
    'DSCR > 1.0x \u2192 Borrower has a cushion \u2014 e.g., 1.32x means $1.32 earned per $1.00 owed',
    'DSCR < 1.0x \u2192 Borrower does not earn enough to cover payments (red flag)',
    'Higher is better \u2014 banks typically want to see 1.2x or above',
]
add_bullet_list(slide, Inches(0.8), Inches(4.5), Inches(6), Inches(2.5), meanings,
                font_size=13, font_color=DARK_GRAY)

# Right side: why it matters
why_box = add_rounded_rect(slide, Inches(7.5), Inches(4.0), Inches(5.2), Inches(3),
                            RGBColor(0xFF, 0xF7, 0xED), ORANGE)
add_text_box(slide, Inches(7.8), Inches(4.15), Inches(4.6), Inches(0.4),
             'Why DSCR Matters', font_size=15, font_color=ORANGE, bold=True)
add_bullet_list(slide, Inches(7.8), Inches(4.55), Inches(4.6), Inches(2.3), [
    'Core credit risk metric for lending',
    'Required for regulatory reporting',
    'Used by rating agencies (S&P, Moody\'s)',
    'Drives covenant compliance decisions',
    'Monitored quarterly across all portfolios',
], font_size=12, font_color=DARK_GRAY)

add_footer(slide, 3)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4: TWO VARIANTS - CRE vs C&I
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_step_badge(slide, Inches(0.8), Inches(0.5), '1', 'Metric Configuration')
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
             'Same formula engine, different inputs depending on the product type',
             font_size=14, font_color=MED_GRAY)

# CRE Card
cre_card = add_rounded_rect(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(5.2),
                             RGBColor(0xF0, 0xF5, 0xFF), BLUE)
add_text_box(slide, Inches(1.1), Inches(1.85), Inches(4), Inches(0.4),
             '\u25cf  CRE DSCR (NOI)', font_size=18, font_color=BLUE, bold=True)
add_text_box(slide, Inches(1.1), Inches(2.3), Inches(5), Inches(0.3),
             'Commercial Real Estate  \u2022  Multifamily  \u2022  Quarterly Monitoring',
             font_size=11, font_color=MED_GRAY)

add_text_box(slide, Inches(1.1), Inches(2.8), Inches(5), Inches(0.3),
             'NUMERATOR \u2014 Net Operating Income (NOI)', font_size=12, font_color=DARK_GRAY, bold=True)
cre_num_items = [
    '+  Gross Potential Rent            $2,400,000',
    '+  Other Income                          $85,000',
    '\u2013   Vacancy & Credit Loss          $120,000',
    '\u2013   Operating Expenses              $780,000',
]
add_bullet_list(slide, Inches(1.3), Inches(3.1), Inches(5), Inches(1.5), cre_num_items,
                font_size=11, font_color=DARK_GRAY, bullet_char=' ')
# NOI total line
add_shape_rect(slide, Inches(1.3), Inches(4.35), Inches(4.5), Pt(1), MED_GRAY)
add_text_box(slide, Inches(1.3), Inches(4.4), Inches(4.5), Inches(0.3),
             'NOI = $1,585,000', font_size=12, font_color=BLUE, bold=True, font_name='Consolas')

add_text_box(slide, Inches(1.1), Inches(4.9), Inches(5), Inches(0.3),
             'DENOMINATOR \u2014 Senior Debt Service', font_size=12, font_color=DARK_GRAY, bold=True)
cre_den_items = [
    '+  Senior Interest                   $720,000',
    '+  Senior Principal                  $480,000',
]
add_bullet_list(slide, Inches(1.3), Inches(5.2), Inches(5), Inches(0.8), cre_den_items,
                font_size=11, font_color=DARK_GRAY, bullet_char=' ')
# Senior DS total line
add_shape_rect(slide, Inches(1.3), Inches(5.8), Inches(4.5), Pt(1), MED_GRAY)
add_text_box(slide, Inches(1.3), Inches(5.85), Inches(4.5), Inches(0.3),
             'Senior DS = $1,200,000', font_size=12, font_color=BLUE, bold=True, font_name='Consolas')

# CRE Result
result_box = add_rounded_rect(slide, Inches(1.1), Inches(6.25), Inches(5.2), Inches(0.55),
                               RGBColor(0xE0, 0xEF, 0xFF), BLUE)
add_text_box(slide, Inches(1.3), Inches(6.3), Inches(4.8), Inches(0.45),
             'NOI \u00f7 Senior DS  =  1.32x', font_size=20, font_color=BLUE, bold=True,
             alignment=PP_ALIGN.CENTER, font_name='Consolas')

# C&I Card
ci_card = add_rounded_rect(slide, Inches(6.9), Inches(1.7), Inches(5.8), Inches(5.2),
                            RGBColor(0xF5, 0xF0, 0xFF), PURPLE)
add_text_box(slide, Inches(7.2), Inches(1.85), Inches(4), Inches(0.4),
             '\u25cf  C&I DSCR (EBITDA)', font_size=18, font_color=PURPLE, bold=True)
add_text_box(slide, Inches(7.2), Inches(2.3), Inches(5), Inches(0.3),
             'Corporate / C&I  \u2022  Middle Market  \u2022  Annual Review',
             font_size=11, font_color=MED_GRAY)

add_text_box(slide, Inches(7.2), Inches(2.8), Inches(5), Inches(0.3),
             'NUMERATOR \u2014 EBITDA', font_size=12, font_color=DARK_GRAY, bold=True)
ci_num_items = [
    '+  Net Income                         $3,200,000',
    '+  Interest Expense (add-back)    $890,000',
    '+  Tax Provision (add-back)        $1,100,000',
    '+  Depreciation & Amort.            $450,000',
]
add_bullet_list(slide, Inches(7.4), Inches(3.1), Inches(5), Inches(1.5), ci_num_items,
                font_size=11, font_color=DARK_GRAY, bullet_char=' ')
# EBITDA total line
add_shape_rect(slide, Inches(7.4), Inches(4.35), Inches(4.5), Pt(1), MED_GRAY)
add_text_box(slide, Inches(7.4), Inches(4.4), Inches(4.5), Inches(0.3),
             'EBITDA = $5,640,000', font_size=12, font_color=PURPLE, bold=True, font_name='Consolas')

add_text_box(slide, Inches(7.2), Inches(4.9), Inches(5), Inches(0.3),
             'DENOMINATOR \u2014 Global Debt Service', font_size=12, font_color=DARK_GRAY, bold=True)
ci_den_items = [
    '+  Senior Interest                    $720,000',
    '+  Senior Principal                   $480,000',
    '+  Mezz / Sub Debt P&I             $180,000',
]
add_bullet_list(slide, Inches(7.4), Inches(5.2), Inches(5), Inches(1), ci_den_items,
                font_size=11, font_color=DARK_GRAY, bullet_char=' ')
# Global DS total line
add_shape_rect(slide, Inches(7.4), Inches(5.95), Inches(4.5), Pt(1), MED_GRAY)
add_text_box(slide, Inches(7.4), Inches(6.0), Inches(4.5), Inches(0.3),
             'Global DS = $1,380,000', font_size=12, font_color=PURPLE, bold=True, font_name='Consolas')

# C&I Result
result_box2 = add_rounded_rect(slide, Inches(7.2), Inches(6.25), Inches(5.2), Inches(0.55),
                                RGBColor(0xF0, 0xE8, 0xFF), PURPLE)
add_text_box(slide, Inches(7.4), Inches(6.3), Inches(4.8), Inches(0.45),
             'EBITDA \u00f7 Global DS  =  4.09x', font_size=20, font_color=PURPLE, bold=True,
             alignment=PP_ALIGN.CENTER, font_name='Consolas')

add_footer(slide, 4)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5: STEP 1 SCREENSHOT
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_step_badge(slide, Inches(0.8), Inches(0.5), '1', 'Metric Configuration \u2014 Live View')
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.4),
             'Both CRE and C&I variants shown side by side with full formula breakdowns',
             font_size=13, font_color=MED_GRAY)

# Screenshot with border frame
border = add_rounded_rect(slide, Inches(0.7), Inches(1.45), Inches(11.9), Inches(5.55),
                           RGBColor(0xE8, 0xE8, 0xEE), RGBColor(0xCC, 0xCC, 0xDD))
add_screenshot(slide, f'{SCREENSHOTS}/01_header_step1.png',
               Inches(0.8), Inches(1.5), Inches(11.7))

add_footer(slide, 5)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6: STEP 2 - DIMENSIONAL ANCHORS
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_step_badge(slide, Inches(0.8), Inches(0.5), '2', 'Dimensional Anchors (L1 Reference Data)')
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
             'Reference tables that identify the "who" and "what" \u2014 unchanged across variants',
             font_size=14, font_color=MED_GRAY)

# Two column layout: text + screenshot
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(4.5), Inches(0.4),
             'Core Master Tables', font_size=15, font_color=DARK_GRAY, bold=True)
add_bullet_list(slide, Inches(0.8), Inches(2.0), Inches(4.5), Inches(2), [
    'facility_master \u2014 The individual loan or property',
    'counterparty \u2014 The borrower (company or person)',
    'credit_agreement_master \u2014 The legal contract',
], font_size=12, font_color=DARK_GRAY)

add_text_box(slide, Inches(0.8), Inches(3.5), Inches(4.5), Inches(0.4),
             'Dimensional & Classification Tables', font_size=15, font_color=DARK_GRAY, bold=True)
add_bullet_list(slide, Inches(0.8), Inches(3.9), Inches(4.5), Inches(2.5), [
    'enterprise_business_taxonomy \u2014 LoB hierarchy for rollup',
    'enterprise_product_taxonomy \u2014 CRE, C&I, PF, FF classification',
    'scenario_dim \u2014 Stress testing scenarios',
    'currency_dim \u2014 Normalizes amounts across currencies',
    'date_dim \u2014 Calendar and fiscal periods',
    'metric_definition_dim \u2014 Defines DSCR as a metric',
], font_size=12, font_color=DARK_GRAY)

# Callout box
callout = add_rounded_rect(slide, Inches(0.8), Inches(6.0), Inches(4.5), Inches(0.8),
                            RGBColor(0xFE, 0xF3, 0xC7), ORANGE)
add_text_box(slide, Inches(1.0), Inches(6.1), Inches(4.1), Inches(0.6),
             '9 dimensional tables anchor every DSCR calculation. '
             'These provide the hierarchy that enables rollup from facility to LoB.',
             font_size=11, font_color=DARK_GRAY)

# Screenshot on right
add_screenshot(slide, f'{SCREENSHOTS}/02_step2_l1_ref.png',
               Inches(5.8), Inches(1.5), Inches(7), with_border=True)

add_footer(slide, 6)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7: STEP 3 - SOURCE DATA TABLES
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_step_badge(slide, Inches(0.8), Inches(0.5), '3', 'Source Data Tables (L2 Snapshot)')
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
             'Two tables serve both variants \u2014 different fields activate per product',
             font_size=14, font_color=MED_GRAY)

# Left: Table descriptions
add_text_box(slide, Inches(0.8), Inches(1.7), Inches(5), Inches(0.4),
             'L2.facility_financial_snapshot', font_size=14, font_color=TEAL, bold=True,
             font_name='Consolas')
add_bullet_list(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(1.5), [
    'Quarterly financial data for each facility',
    'CRE fields: revenue_amt, operating_expense_amt, noi_amt',
    'C&I fields: ebitda_amt, interest_expense_amt',
    'Shared: principal_payment_amt, total_debt_service_amt',
], font_size=12, font_color=DARK_GRAY)

add_text_box(slide, Inches(0.8), Inches(3.6), Inches(5), Inches(0.4),
             'L2.cash_flow', font_size=14, font_color=TEAL, bold=True,
             font_name='Consolas')
add_bullet_list(slide, Inches(0.8), Inches(4.0), Inches(5.5), Inches(1.3), [
    'Individual cash flow events (interest & principal)',
    'Both CRE and C&I variants use this table',
    'Filters by cash_flow_type: interest vs. principal',
    'Joins to facility_master and counterparty via FK',
], font_size=12, font_color=DARK_GRAY)

# Highlight box
hl = add_rounded_rect(slide, Inches(0.8), Inches(5.3), Inches(5.5), Inches(0.7),
                       RGBColor(0xEC, 0xFD, 0xF5), GREEN)
add_text_box(slide, Inches(1.0), Inches(5.4), Inches(5.1), Inches(0.5),
             'Each field is tagged CRE, C&I, or Both \u2014 the platform activates '
             'only the relevant fields based on the selected variant.',
             font_size=11, font_color=DARK_GRAY)

# Screenshot
add_screenshot(slide, f'{SCREENSHOTS}/03_step3_l2_snapshot.png',
               Inches(6.5), Inches(1.3), Inches(6.3), with_border=True)

add_footer(slide, 7)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8: STEP 4 - CALCULATION
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_step_badge(slide, Inches(0.8), Inches(0.5), '4', 'Facility-Level Calculation')
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
             'The DSCR formula is applied at the individual facility (loan) level',
             font_size=14, font_color=MED_GRAY)

# T2 Authority explanation
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.4),
             'T2 Calculation Authority', font_size=16, font_color=DARK_GRAY, bold=True)
add_bullet_list(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(2), [
    'The bank (GSIB) sends their own DSCR value',
    'The platform independently recalculates from raw data',
    'If values differ beyond tolerance, a reconciliation flag is raised',
    'This is the only level where the bank provides DSCR',
    'All levels above facility are always calculated by the platform (T3)',
], font_size=12, font_color=DARK_GRAY)

# Two result boxes
# CRE
cre_res = add_rounded_rect(slide, Inches(0.8), Inches(4.2), Inches(5.5), Inches(1.5),
                            RGBColor(0xF0, 0xF5, 0xFF), BLUE)
add_text_box(slide, Inches(1.1), Inches(4.3), Inches(3), Inches(0.3),
             'CRE DSCR (NOI)', font_size=14, font_color=BLUE, bold=True)
add_text_box(slide, Inches(1.1), Inches(4.7), Inches(5), Inches(0.8),
             'NOI: $1,585,000\nSenior DS: $1,200,000\n\nDSCR = 1.32x',
             font_size=13, font_color=DARK_GRAY, font_name='Consolas')

# C&I
ci_res = add_rounded_rect(slide, Inches(0.8), Inches(5.9), Inches(5.5), Inches(1.5),
                           RGBColor(0xF5, 0xF0, 0xFF), PURPLE)
add_text_box(slide, Inches(1.1), Inches(6.0), Inches(3), Inches(0.3),
             'C&I DSCR (EBITDA)', font_size=14, font_color=PURPLE, bold=True)
add_text_box(slide, Inches(1.1), Inches(6.4), Inches(5), Inches(0.8),
             'EBITDA: $5,640,000\nGlobal DS: $1,380,000\n\nDSCR = 4.09x',
             font_size=13, font_color=DARK_GRAY, font_name='Consolas')

# Screenshot
add_screenshot(slide, f'{SCREENSHOTS}/04_step4_calculation.png',
               Inches(6.8), Inches(1.3), Inches(6), with_border=True)

add_footer(slide, 8)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9: THE FOUNDATIONAL RULE
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
             'The Golden Rule of DSCR Aggregation', font_size=32, font_color=BLACK, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(3), RED)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.7),
             'At every level above facility, DSCR must be re-derived from raw components. '
             'You never take facility DSCRs and average them.',
             font_size=16, font_color=DARK_GRAY)

# Wrong vs Right
# WRONG
wrong_box = add_rounded_rect(slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(2.5),
                              RGBColor(0xFE, 0xF2, 0xF2), RED)
add_text_box(slide, Inches(1.1), Inches(2.65), Inches(3), Inches(0.4),
             '\u2717  WRONG \u2014 Simple Average', font_size=16, font_color=RED, bold=True)
add_text_box(slide, Inches(1.1), Inches(3.1), Inches(5), Inches(0.5),
             'avg(DSCR\u2081, DSCR\u2082, \u2026, DSCR\u2099)', font_size=18, font_color=RED,
             font_name='Consolas', alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.1), Inches(3.7), Inches(5), Inches(1),
             'A $100M facility with DSCR 0.8x and a $1M facility with DSCR 5.0x '
             'would average to 2.9x \u2014 but the portfolio is overwhelmingly at risk!',
             font_size=12, font_color=DARK_GRAY)

# RIGHT
right_box = add_rounded_rect(slide, Inches(6.9), Inches(2.5), Inches(5.8), Inches(2.5),
                              RGBColor(0xEC, 0xFD, 0xF5), GREEN)
add_text_box(slide, Inches(7.2), Inches(2.65), Inches(4), Inches(0.4),
             '\u2713  CORRECT \u2014 Pool & Divide', font_size=16, font_color=GREEN, bold=True)
add_text_box(slide, Inches(7.2), Inches(3.1), Inches(5), Inches(0.5),
             '\u03a3 Cash Flow  \u00f7  \u03a3 Debt Service', font_size=18, font_color=GREEN,
             font_name='Consolas', alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(7.2), Inches(3.7), Inches(5), Inches(1),
             'Sum all the numerators (cash flows) across facilities, sum all the '
             'denominators (debt service), then divide once. This properly weights '
             'by exposure size.',
             font_size=12, font_color=DARK_GRAY)

# Analogy box
analogy = add_rounded_rect(slide, Inches(0.8), Inches(5.3), Inches(5.5), Inches(1.8),
                            RGBColor(0xFE, 0xF3, 0xC7), ORANGE)
add_text_box(slide, Inches(1.1), Inches(5.4), Inches(5), Inches(0.3),
             'Plain English Analogy', font_size=14, font_color=ORANGE, bold=True)
add_text_box(slide, Inches(1.1), Inches(5.8), Inches(5), Inches(1.2),
             'Think of it like a school\'s GPA: you don\'t average GPAs from each classroom. '
             'You add all grade points and credit hours from every student, then divide once. '
             'Rating agencies like S&P calculate pooled DSCR the same way.',
             font_size=12, font_color=DARK_GRAY, line_spacing=1.4)

# Screenshot on right bottom
add_screenshot(slide, f'{SCREENSHOTS}/06_foundational_rule.png',
               Inches(6.9), Inches(5.2), Inches(5.8), with_border=True)

add_footer(slide, 9)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10: STEP 5 - L3 OUTPUT TABLES
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_step_badge(slide, Inches(0.8), Inches(0.5), '5a', 'L3 Output Tables')
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
             'Calculated results are stored in L3 tables for consumption',
             font_size=14, font_color=MED_GRAY)

# Table list
tables = [
    ('L3.metric_value_fact', 'Primary', 'All DSCR variants at every aggregation level (facility, counterparty, desk, portfolio, LoB)'),
    ('L3.lob_credit_quality_summary', '', 'LoB-level credit quality including DSCR value'),
    ('L3.lob_risk_ratio_summary', '', 'LoB-level risk ratios \u2014 DSCR alongside FCCR, LTV, capital adequacy'),
    ('L3.risk_appetite_metric_state', '', 'DSCR vs. risk appetite limits with RAG status and velocity'),
    ('L3.facility_detail_snapshot', '', 'Facility-level analytics for drill-down pop-ups'),
]

for i, (name, badge, desc) in enumerate(tables):
    y = Inches(1.6) + i * Inches(0.85)
    color = GREEN if badge == 'Primary' else RGBColor(0xF0, 0xF0, 0xF5)
    border = GREEN if badge == 'Primary' else RGBColor(0xDD, 0xDD, 0xDD)
    card = add_rounded_rect(slide, Inches(0.8), y, Inches(6), Inches(0.7), color, border)
    add_text_box(slide, Inches(1.0), y + Inches(0.05), Inches(4), Inches(0.3),
                 name, font_size=13, font_color=WHITE if badge == 'Primary' else DARK_GRAY,
                 bold=True, font_name='Consolas')
    add_text_box(slide, Inches(1.0), y + Inches(0.35), Inches(5.5), Inches(0.3),
                 desc, font_size=10, font_color=WHITE if badge == 'Primary' else MED_GRAY)

# Screenshot
add_screenshot(slide, f'{SCREENSHOTS}/05_step5_l3_output.png',
               Inches(7.2), Inches(1.3), Inches(5.8), with_border=True)

add_footer(slide, 10)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11: ROLLUP HIERARCHY
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_step_badge(slide, Inches(0.8), Inches(0.5), '5b', 'Rollup Hierarchy')
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
             'DSCR aggregates through 5 levels, each using the "pool and divide" method',
             font_size=14, font_color=MED_GRAY)

# Pyramid levels
levels = [
    ('Facility', 'Direct Calculation', 'T2 \u2014 Source + Validate', 'Formula applied directly to each loan', BLUE, Inches(5.8)),
    ('Counterparty', 'Pooled Ratio', 'T3 \u2014 Always Calculate', 'Sum facility numerators & denominators per borrower', TEAL, Inches(5.0)),
    ('Desk', 'Product-Segmented Pooled', 'T3 \u2014 Always Calculate', 'Pool by product type within each desk', GREEN, Inches(4.2)),
    ('Portfolio', 'Exposure-Weighted Avg', 'T3 \u2014 Always Calculate', 'Weighted average using committed exposure', PURPLE, Inches(3.4)),
    ('Line of Business', 'Exposure-Weighted Avg', 'T3 \u2014 Always Calculate', 'Highest level for directional monitoring', RGBColor(0xEC, 0x48, 0x99), Inches(2.6)),
]

for i, (name, method, authority, desc, color, bar_w) in enumerate(levels):
    y = Inches(1.7) + i * Inches(1.1)
    x_center = Inches(3.2)  # Center of left half (0 to 6.3)
    x_start = x_center - bar_w / 2

    # Level bar (pyramid effect)
    bar = add_rounded_rect(slide, x_start, y, bar_w, Inches(0.85), RGBColor(0xFA, 0xFA, 0xFF), color)

    # Level name
    add_text_box(slide, x_start + Inches(0.2), y + Inches(0.05), Inches(2), Inches(0.3),
                 name, font_size=14, font_color=color, bold=True)
    # Method tag
    add_text_box(slide, x_start + Inches(0.2), y + Inches(0.32), bar_w - Inches(0.4), Inches(0.2),
                 method, font_size=10, font_color=MED_GRAY, bold=True)
    # Description
    add_text_box(slide, x_start + Inches(0.2), y + Inches(0.52), bar_w - Inches(0.4), Inches(0.3),
                 desc, font_size=10, font_color=DARK_GRAY)

# Arrow indicators between levels
for i in range(4):
    y = Inches(2.55) + i * Inches(1.1)
    add_text_box(slide, Inches(3.0), y, Inches(0.5), Inches(0.3),
                 '\u25b2', font_size=14, font_color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# Screenshot on right
add_screenshot(slide, f'{SCREENSHOTS}/07_rollup_hierarchy.png',
               Inches(6.3), Inches(1.3), Inches(6.5), with_border=True)

add_footer(slide, 11)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12: STEP 6 - DASHBOARD CONNECTION
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_step_badge(slide, Inches(0.8), Inches(0.5), '6', 'Self-Service Dashboard Connection')
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
             'Pick the metric, pick the dimension, build \u2014 no SQL needed',
             font_size=14, font_color=MED_GRAY)

# 3-step process
process_steps = [
    ('1', 'Select Variant', 'Choose CRE DSCR (NOI)\nor C&I DSCR (EBITDA)', BLUE),
    ('2', 'Select Dimension', 'Facility, Counterparty,\nDesk, Portfolio, or LoB', GREEN),
    ('3', 'Dashboard Output', 'Platform auto-resolves\ntable joins \u2014 zero SQL', ORANGE),
]

for i, (num, title, desc, color) in enumerate(process_steps):
    x = Inches(0.8) + i * Inches(2.3)
    y = Inches(1.7)

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = circle.text_frame.paragraphs[0].add_run()
    run.text = num
    run.font.size = Pt(18)
    run.font.color.rgb = WHITE
    run.font.bold = True

    add_text_box(slide, x, y + Inches(0.7), Inches(2), Inches(0.3),
                 title, font_size=13, font_color=DARK_GRAY, bold=True)
    add_text_box(slide, x, y + Inches(1.0), Inches(2), Inches(0.7),
                 desc, font_size=11, font_color=MED_GRAY)

# Arrow connectors
for i in range(2):
    x = Inches(2.5) + i * Inches(2.3)
    add_text_box(slide, x, Inches(1.9), Inches(0.5), Inches(0.3),
                 '\u25b6', font_size=16, font_color=MED_GRAY)

# Key benefits box
benefits_box = add_rounded_rect(slide, Inches(0.8), Inches(3.8), Inches(5.5), Inches(1.5),
                                 RGBColor(0xEC, 0xFD, 0xF5), GREEN)
add_text_box(slide, Inches(1.1), Inches(3.9), Inches(4), Inches(0.3),
             'Key Benefits', font_size=14, font_color=GREEN, bold=True)
add_bullet_list(slide, Inches(1.1), Inches(4.2), Inches(5), Inches(1), [
    'No SQL required \u2014 platform handles all joins',
    'Automatic lineage from definition to dashboard',
    'Consistent aggregation at every level',
    'Full audit trail for regulatory compliance',
], font_size=12, font_color=DARK_GRAY)

# Screenshot
add_screenshot(slide, f'{SCREENSHOTS}/08_step6_dashboard.png',
               Inches(6.8), Inches(1.3), Inches(6), with_border=True)

add_footer(slide, 12)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13: CALCULATION AUTHORITY TIERS
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
             'Calculation Authority Tiers', font_size=32, font_color=BLACK, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(2.5), ORANGE)
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(10), Inches(0.5),
             'Who calculates the number at each stage, and how it gets verified',
             font_size=14, font_color=MED_GRAY)

# T1 Box
t1 = add_rounded_rect(slide, Inches(0.8), Inches(2.2), Inches(3.7), Inches(4.5),
                       RGBColor(0xF0, 0xF5, 0xFF), BLUE)
add_text_box(slide, Inches(1.1), Inches(2.35), Inches(3), Inches(0.4),
             'T1 \u2014 Always Source', font_size=18, font_color=BLUE, bold=True)
add_text_box(slide, Inches(1.1), Inches(2.85), Inches(3.2), Inches(0.3),
             'Where it applies:', font_size=12, font_color=DARK_GRAY, bold=True)
add_text_box(slide, Inches(1.1), Inches(3.15), Inches(3.2), Inches(0.3),
             'L1 Reference Data, L2 Source Tables', font_size=12, font_color=MED_GRAY)
add_bullet_list(slide, Inches(1.1), Inches(3.6), Inches(3.2), Inches(2.5), [
    'Data is sourced directly from bank systems',
    'No calculation by the platform',
    'Platform stores it as-is for downstream use',
    'Examples: facility_master, counterparty, cash_flow',
], font_size=11, font_color=DARK_GRAY)

# T2 Box
t2 = add_rounded_rect(slide, Inches(4.8), Inches(2.2), Inches(3.7), Inches(4.5),
                       RGBColor(0xEC, 0xFD, 0xF5), GREEN)
add_text_box(slide, Inches(5.1), Inches(2.35), Inches(3), Inches(0.4),
             'T2 \u2014 Source + Validate', font_size=18, font_color=GREEN, bold=True)
add_text_box(slide, Inches(5.1), Inches(2.85), Inches(3.2), Inches(0.3),
             'Where it applies:', font_size=12, font_color=DARK_GRAY, bold=True)
add_text_box(slide, Inches(5.1), Inches(3.15), Inches(3.2), Inches(0.3),
             'Facility-level DSCR only', font_size=12, font_color=MED_GRAY)
add_bullet_list(slide, Inches(5.1), Inches(3.6), Inches(3.2), Inches(2.5), [
    'Bank sends their DSCR value',
    'Platform recalculates independently',
    'Values are compared',
    'Differences trigger reconciliation flags',
    'This is the only level banks provide DSCR',
], font_size=11, font_color=DARK_GRAY)

# T3 Box
t3 = add_rounded_rect(slide, Inches(8.8), Inches(2.2), Inches(3.9), Inches(4.5),
                       RGBColor(0xF5, 0xF0, 0xFF), PURPLE)
add_text_box(slide, Inches(9.1), Inches(2.35), Inches(3.2), Inches(0.4),
             'T3 \u2014 Always Calculate', font_size=18, font_color=PURPLE, bold=True)
add_text_box(slide, Inches(9.1), Inches(2.85), Inches(3.2), Inches(0.3),
             'Where it applies:', font_size=12, font_color=DARK_GRAY, bold=True)
add_text_box(slide, Inches(9.1), Inches(3.15), Inches(3.2), Inches(0.3),
             'Counterparty, Desk, Portfolio, LoB', font_size=12, font_color=MED_GRAY)
add_bullet_list(slide, Inches(9.1), Inches(3.6), Inches(3.2), Inches(2.5), [
    'Platform always calculates \u2014 bank never provides',
    'Uses "pool and divide" method',
    'Sums raw numerators and denominators',
    'Then divides once at the target level',
    'Ensures proper exposure weighting',
], font_size=11, font_color=DARK_GRAY)

add_footer(slide, 13)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14: SUMMARY / KEY TAKEAWAYS
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

# Left accent bar
add_shape_rect(slide, Inches(0), Inches(0), Inches(0.4), SLIDE_HEIGHT, ORANGE)

add_text_box(slide, Inches(1.2), Inches(0.8), Inches(10), Inches(0.7),
             'Key Takeaways', font_size=36, font_color=BLACK, bold=True)
add_accent_bar(slide, Inches(1.2), Inches(1.55), Inches(2.5), ORANGE)

takeaways = [
    ('DSCR = Cash Flow \u00f7 Debt Service',
     'Measures how many dollars of income a borrower earns for every dollar of loan payments owed.'),
    ('Two Variants, One Engine',
     'CRE uses Net Operating Income (property income). C&I uses EBITDA (corporate earnings). Same formula structure, different inputs.'),
    ('6-Layer Data Lineage',
     'From user definition \u2192 L1 reference data \u2192 L2 source snapshots \u2192 facility calculation \u2192 L3 storage \u2192 dashboard.'),
    ('Never Average Ratios',
     'Always re-derive by summing raw components (numerators and denominators) and dividing once. Averaging pre-computed ratios gives misleading results.'),
    ('5-Level Rollup Hierarchy',
     'Facility \u2192 Counterparty \u2192 Desk \u2192 Portfolio \u2192 Line of Business. Each level uses pool-and-divide aggregation.'),
    ('Self-Service Dashboards',
     'Pick a variant and dimension \u2014 the platform auto-resolves all table joins. No SQL required.'),
]

colors = [BLUE, PURPLE, TEAL, RED, GREEN, ORANGE]
for i, (title, desc) in enumerate(takeaways):
    col = i % 2
    row = i // 2
    x = Inches(1.2) + col * Inches(5.9)
    y = Inches(2.1) + row * Inches(1.7)

    # Colored dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y + Inches(0.08), Inches(0.22), Inches(0.22))
    dot.fill.solid()
    dot.fill.fore_color.rgb = colors[i]
    dot.line.fill.background()

    add_text_box(slide, x + Inches(0.35), y, Inches(5.2), Inches(0.35),
                 title, font_size=15, font_color=DARK_GRAY, bold=True)
    add_text_box(slide, x + Inches(0.35), y + Inches(0.4), Inches(5.2), Inches(1),
                 desc, font_size=12, font_color=MED_GRAY, line_spacing=1.3)

add_footer(slide, 14)


# ═══════════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════════
prs.save(OUTPUT)
print(f'Saved to: {OUTPUT}')
print(f'Total slides: {len(prs.slides)}')
